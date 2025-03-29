import openai
import os
import pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from typing import List, Dict
import argparse
import hashlib
from datetime import datetime
from openai import AsyncOpenAI
import asyncio
import cProfile
import re

# --- Constants and Configurations ---

LOG_FILE = "translation_log.txt"
SUPPORTED_EXTENSIONS = ('.pptx',)

# --- Helper Functions ---
def log_message(message: str, level: str = "INFO") -> None:
    """Logs messages to the console and a log file with timestamps."""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_entry = f"[{timestamp}] [{level}] {message}\n"
    print(log_entry, end="")
    with open(LOG_FILE, "a", encoding='utf-8') as logfile:
        logfile.write(log_entry)

def get_cache_filename(presentation_path: str, target_language: str) -> str:
    """Generates a unique cache filename based on the presentation and language."""
    base_name = os.path.splitext(os.path.basename(presentation_path))[0]
    path_hash = hashlib.sha256(presentation_path.encode('utf-8')).hexdigest()[:8]  # Short hash
    return f"translation_cache_{base_name}_{target_language}_{path_hash}.json"

def load_cache(cache_file: str) -> Dict:
    """Loads the translation cache from the specified file (if it exists)."""
    try:
        with open(cache_file, "r", encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        return {}

def save_cache(cache: Dict, cache_file: str) -> None:
    """Saves the translation cache to the specified file."""
    with open(cache_file, "w", encoding='utf-8') as f:
        json.dump(cache, f, indent=4, ensure_ascii=False)

def generate_prompt_hash(prompt: str) -> str:
    """Generates a SHA-256 hash of the prompt for use as a cache key."""
    return hashlib.sha256(prompt.encode('utf-8')).hexdigest()

async def translate_text_with_openai(prompt: str, target_language: str, cache: Dict, max_retries: int = 3) -> str:
    """Translates text using the OpenAI API, with caching and retries."""
    prompt_hash = generate_prompt_hash(prompt)
    if prompt_hash in cache:
        log_message(f"Using cached translation for: {prompt[:50]}...", level="CACHE_HIT")
        return cache[prompt_hash]

    log_message(f"Translating: {prompt[:50]}... to {target_language}", level="API_CALL")

    schema_instruction = (
        "Return the translation as a JSON object exactly as follows: "
        "{\"translated\": \"<translated text>\"}"
    )

    system_instruction = (
        f"You are a helpful assistant that translates text to {target_language}. "
        "Maintain the original meaning as closely as possible. "
        f"Adjust the tone of the translation to be appropriate for professional presentations in the target language ({target_language}). "
        "The translated text should be approximately the same character length as the original text (within a 5% margin). Do not translate things like emails, phone numbers, or other non-text content. "
        f"{schema_instruction}"
    )

    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OpenAI API key not found. Please set the OPENAI_API_KEY environment variable.")

    client = AsyncOpenAI(api_key=api_key, timeout=30.0)

    for attempt in range(max_retries):
        try:
            response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.2,
                max_tokens=4096,
                timeout=30,
                response_format={"type": "json_object"}
            )
            result = response.choices[0].message.content.strip()

            for parse_attempt in range(max_retries):
                try:
                    parsed = json.loads(result)
                    translated_text = parsed["translated"]
                    cache[prompt_hash] = translated_text
                    return translated_text  # Return immediately after saving to cache
                except (json.JSONDecodeError, KeyError) as e:
                    log_message(f"Error parsing response (attempt {parse_attempt+1}/{max_retries}): {str(e)}", level="ERROR")
                    if parse_attempt == max_retries - 1:
                        log_message(f"Failed to parse response after {max_retries} attempts. Returning original prompt.", level="ERROR")
                        return prompt
                except Exception as e:
                    log_message(f"Unexpected error during parsing: {e}", level="ERROR")
                    return prompt

        except Exception as e:
            log_message(f"OpenAI API Error (Attempt {attempt + 1}/{max_retries}): {str(e)}", level="ERROR")
            if attempt == max_retries - 1:
                return prompt
    return prompt


def extract_text_from_presentation(presentation_path: str) -> List[Dict]:
    """Extracts text and context from a PowerPoint presentation."""
    try:
        prs = Presentation(presentation_path)
        text_data = []

        def _extract_text_from_shape(shape, slide_number, base_shape_id):
            """Helper function to recursively extract text from shapes and groups."""
            nonlocal text_data
            shape_index = base_shape_id # Use the passed-in index/id

            # Handle Groups Recursively
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                log_message(f"Processing Group Shape: {shape.name} ({base_shape_id})", level="DEBUG")
                for i, sub_shape in enumerate(shape.shapes):
                    sub_shape_id = f"{base_shape_id}_sub{i}"
                    _extract_text_from_shape(sub_shape, slide_number, sub_shape_id)
                return # Don't process the group container itself further

            # Handle Text Frames
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text.strip(): # Check if frame has text before iterating
                    log_message(f"Processing Text Frame in Shape: {shape.name} ({base_shape_id})", level="DEBUG")
                    for para_idx, paragraph in enumerate(text_frame.paragraphs):
                        for run_idx, run in enumerate(paragraph.runs):
                            if run.text.strip():
                                shape_type = "UNKNOWN"
                                try:
                                    # Checking title requires access to the slide object, difficult here
                                    # Simplified type detection
                                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                                        shape_type = "PLACEHOLDER"
                                    else:
                                        shape_type = "BODY" # Default assumption
                                except AttributeError:
                                    shape_type = "BODY" # Fallback

                                text_data.append({
                                    "slide_number": slide_number,
                                    "shape_type": shape_type,
                                    "text": run.text,
                                    "shape_id": f"{base_shape_id}_p{para_idx}_r{run_idx}",
                                    "original_shape_id": base_shape_id,
                                    "paragraph_index": para_idx,
                                    "run_index": run_idx,
                                })

            # Handle Tables
            elif shape.has_table:
                log_message(f"Processing Table Shape: {shape.name} ({base_shape_id})", level="DEBUG")
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if cell_text:
                            text_data.append({
                                "slide_number": slide_number,
                                "shape_type": "TABLE",
                                "text": cell_text,
                                "shape_id": f"{base_shape_id}_row{row_idx}_col{col_idx}",
                                "original_shape_id": base_shape_id,
                                "row_index": row_idx,
                                "col_index": col_idx,
                            })
            # else: # Optional: Log shapes that are neither text, table, nor group
                # if shape.shape_type != MSO_SHAPE_TYPE.GROUP: # Already handled
                    # log_message(f"Skipping non-text/table/group shape: {shape.name} ({shape.shape_type}) ({base_shape_id})", level="DEBUG")

        # --- Main Extraction Loop ---
        for slide_number, slide in enumerate(prs.slides, start=1):
            log_message(f"Processing Slide {slide_number}", level="DEBUG")
            # Extract from shapes on the slide
            for shape_index, shape in enumerate(slide.shapes):
                shape_id = f"slide{slide_number}_shape{shape_index}"
                _extract_text_from_shape(shape, slide_number, shape_id)

            # Extract from speaker notes
            if slide.has_notes_slide:
                log_message(f"Processing Notes for Slide {slide_number}", level="DEBUG")
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    for para_idx, paragraph in enumerate(notes_tf.paragraphs):
                        for run_idx, run in enumerate(paragraph.runs):
                             if run.text.strip():
                                text_data.append({
                                    "slide_number": slide_number,
                                    "shape_type": "NOTES",
                                    "text": run.text,
                                    "shape_id": f"slide{slide_number}_notes_p{para_idx}_r{run_idx}",
                                    "original_shape_id": f"slide{slide_number}_notes",
                                    "paragraph_index": para_idx,
                                    "run_index": run_idx,
                                })

        return text_data

    except Exception as e:
        log_message(f"Error extracting text from {presentation_path}: {e}", level="ERROR")
        return []


async def batch_translate_texts_with_openai(text_entries: List[Dict], target_language: str, cache: Dict, max_retries: int = 3, batch_size: int = 10) -> None:
    """Batch translates multiple texts using the OpenAI API with structured JSON output."""
    texts_to_translate = []
    for entry in text_entries:
        prompt_hash = generate_prompt_hash(entry["text"])
        if prompt_hash not in cache:
            texts_to_translate.append((prompt_hash, entry["text"]))

    if not texts_to_translate:
        log_message("All translations found in cache.", level="INFO")
        return

    total_batches = (len(texts_to_translate) + batch_size - 1) // batch_size
    log_message(f"Starting batch translation of {len(texts_to_translate)} texts in {total_batches} batches", level="INFO")

    tasks = []
    for batch_num, i in enumerate(range(0, len(texts_to_translate), batch_size), 1):
        batch = texts_to_translate[i:i + batch_size]
        payload = {hash_: text for hash_, text in batch}

        log_message(f"Processing batch {batch_num}/{total_batches} ({len(batch)} texts)", level="INFO")

        schema_instruction = (
            "Return the translations as a JSON object exactly as follows: \n"
            "{\"translations\": {\"<sha256 hash>\": \"<translated text>\"} }"
        )

        system_instruction = (
            f"You are a helpful assistant that translates multiple texts to {target_language}. "
            "Maintain the original meaning as closely as possible. "
            f"Adjust the tone of each translation to be appropriate for professional presentations in the target language ({target_language}). "
            "The translated text for each input should be approximately the same length as the original text (within a 10% margin). "
            f"{schema_instruction}"
        )

        prompt_data = {
            "texts": payload,
            "target_language": target_language,
            "instructions": "Translate each text, maintaining original meaning and formatting."
        }

        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OpenAI API Key not found. Please set the OPENAI_API_KEY environment variable.")

        client = AsyncOpenAI(api_key=api_key, timeout=60.0)

        tasks.append(translate_batch(client, system_instruction, prompt_data, cache, max_retries, batch_num, total_batches))

    await asyncio.gather(*tasks)
    log_message("Batch translation completed", level="INFO")


async def translate_batch(client: AsyncOpenAI, system_instruction:str, prompt_data: Dict, cache: Dict, max_retries: int, batch_num: int, total_batches: int) -> None:
    """Translates a single batch (async).  This is now a separate function."""
    for attempt in range(max_retries):
        try:
            response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": json.dumps(prompt_data)}
                ],
                temperature=0.2,
                max_tokens=4096,
                timeout=60,
                response_format={"type": "json_object"}
            )
            output = response.choices[0].message.content.strip()

            for parse_attempt in range(max_retries):
                try:
                    result = json.loads(output)
                    translations = result.get("translations", {})
                    for hash_, translated_text in translations.items():
                        cache[hash_] = translated_text
                    break  # Success, break out of parsing attempts
                except (json.JSONDecodeError, KeyError) as e:
                    log_message(f"Error parsing batch {batch_num} response (attempt {parse_attempt+1}/{max_retries}): {str(e)}", level="ERROR")
                    if parse_attempt == max_retries - 1:
                        log_message("Failed to parse batch response after all retries", level="ERROR")
                except Exception as e:
                    log_message(f"Unexpected error during parsing: {e}", level="ERROR")
            else:
                continue  # Failed all parsing attempts, retry API call
            break  # Success, break out of API call attempts

        except Exception as e:
            log_message(f"Batch {batch_num} translation error (Attempt {attempt + 1}/{max_retries}): {str(e)}", level="ERROR")
            if attempt == max_retries - 1:
                log_message(f"Max retries reached for batch {batch_num}", level="ERROR")


async def translate_presentation(presentation_path: str, target_language: str, output_path: str) -> None:
    """Translates a PowerPoint presentation and saves the translated version."""

    if not presentation_path.lower().endswith(SUPPORTED_EXTENSIONS):
        log_message(f"Unsupported file format: {presentation_path}. Skipping.", level="WARNING")
        return

    if presentation_path == output_path:
        log_message("Input and output paths are the same. This is not allowed.", level="ERROR")
        return

    try:
        # 1. Extract Text
        # Use the improved extraction logic
        text_data = extract_text_from_presentation(presentation_path)
        if not text_data:
            log_message(f"No text found to translate in {presentation_path}.", level="WARNING")
            return

        # 2. Translate Text (using cache)
        cache_file = get_cache_filename(presentation_path, target_language)
        cache = load_cache(cache_file)
        await batch_translate_texts_with_openai(text_data, target_language, cache)
        save_cache(cache, cache_file)

        # 3. Apply Translations
        # Create a dictionary mapping original text hash to translated text for quick lookup
        translation_map = {generate_prompt_hash(entry["text"]): cache.get(generate_prompt_hash(entry["text"])) for entry in text_data}

        # Load the original presentation *again* to apply changes cleanly
        # Applying to a modified structure can be complex
        prs = Presentation(presentation_path)

        # --- Helper function to find nested shapes (needed for groups) ---
        def find_shape_recursive(parent_shape_collection, target_id_parts):
            if not target_id_parts:
                log_message("find_shape_recursive called with empty target_id_parts", level="WARNING")
                return None
            
            current_part = target_id_parts[0]
            remaining_parts = target_id_parts[1:]

            # Extract index from the current part (e.g., 'shape0', 'sub1')
            try:
                # Find the last sequence of digits in the part
                match = re.search(r'(\d+)$', current_part)
                if not match:
                     log_message(f"Could not find numerical index in shape ID part: {current_part}", level="WARNING")
                     return None
                shape_index = int(match.group(1))

            except (ValueError, IndexError):
                log_message(f"Could not parse index from shape ID part: {current_part}", level="WARNING")
                return None

            if shape_index < len(parent_shape_collection):
                shape = parent_shape_collection[shape_index]
                if not remaining_parts:
                    # Found the target shape
                    return shape
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Recurse into the group
                    return find_shape_recursive(shape.shapes, remaining_parts)
                else:
                    # ID has more parts, but shape is not a group - mismatch
                    log_message(f"ID path continues {remaining_parts}, but shape {shape_index} is not a group.", level="WARNING")
                    return None
            else:
                log_message(f"Shape index {shape_index} out of bounds for collection size {len(parent_shape_collection)}.", level="WARNING")
                return None
        # --- End Helper --- 

        # Apply translations entry by entry
        applied_count = 0
        skipped_count = 0
        for entry in text_data:
            original_text = entry["text"]
            prompt_hash = generate_prompt_hash(original_text)
            translated_text = translation_map.get(prompt_hash)

            if not translated_text or translated_text == original_text:
                # Skip if no translation found or translation is same as original
                skipped_count += 1
                continue

            try:
                slide_number = entry["slide_number"]
                target_slide = prs.slides[slide_number - 1]
                shape_type = entry["shape_type"]
                para_idx = entry["paragraph_index"]
                run_idx = entry["run_index"]
                original_shape_id_str = entry["original_shape_id"] # e.g., "slide1_shape0", "slide3_shape4_sub0", "slide1_notes"

                target_run = None
                shape = None

                if shape_type == "NOTES":
                    if target_slide.has_notes_slide and target_slide.notes_slide.notes_text_frame:
                        notes_tf = target_slide.notes_slide.notes_text_frame
                        if para_idx < len(notes_tf.paragraphs):
                            paragraph = notes_tf.paragraphs[para_idx]
                            if run_idx < len(paragraph.runs):
                                target_run = paragraph.runs[run_idx]
                else:
                    # Parse the shape ID string to find the shape
                    # ID Format examples: slide1_shape0_unnamed0, slide3_shape4_Group17_sub0_Group26_sub1_TextBox13
                    id_parts = original_shape_id_str.split('_')
                    if len(id_parts) < 2 or not id_parts[0].startswith('slide'):
                         log_message(f"Invalid original_shape_id format: {original_shape_id_str}", level="WARNING")
                         continue
                    
                    # Start searching from slide shapes, skip slide part (e.g., 'slide1')
                    shape_path_parts = id_parts[1:] # e.g. ['shape0_unnamed0'], ['shape4_Group17', 'sub0_Group26', 'sub1_TextBox13']
                    
                    # We need to reconstruct the path to the shape including group indices
                    # The `original_shape_id` stored during extraction needs adjustment
                    # Let's use the more detailed `shape_id` instead, which has the full path
                    detailed_shape_id = entry["shape_id"]
                    # e.g. "slide1_shape0_unnamed0_p0_r0", "slide3_shape4_Group17_sub0_Group26_sub1_TextBox13_p0_r0"
                    # e.g. "slide1_shape1_Table_row0_col1_p0_r0"
                    
                    id_parts = detailed_shape_id.split('_')
                    path_to_shape = []
                    for part in id_parts[1:]: # Skip slide number part
                        if part.startswith('p') or part.startswith('r') or part.startswith('row') or part.startswith('col'):
                            break # Stop when we reach paragraph/run/cell parts
                        path_to_shape.append(part) # e.g. ['shape0', 'shape4', 'sub0', 'sub1']
                    
                    if not path_to_shape:
                        log_message(f"Could not extract shape path from ID: {detailed_shape_id}", level="WARNING")
                        continue
                        
                    # Find the shape recursively
                    shape = find_shape_recursive(target_slide.shapes, path_to_shape)

                    if shape:
                        if shape_type == "TABLE":
                            row_idx = entry["row_index"]
                            col_idx = entry["col_index"]
                            if shape.has_table:
                                if row_idx < len(shape.table.rows) and col_idx < len(shape.table.columns):
                                    cell = shape.table.cell(row_idx, col_idx)
                                    if cell.text_frame and para_idx < len(cell.text_frame.paragraphs):
                                        paragraph = cell.text_frame.paragraphs[para_idx]
                                        if run_idx < len(paragraph.runs):
                                            target_run = paragraph.runs[run_idx]
                        elif shape.has_text_frame:
                             if para_idx < len(shape.text_frame.paragraphs):
                                paragraph = shape.text_frame.paragraphs[para_idx]
                                if run_idx < len(paragraph.runs):
                                    target_run = paragraph.runs[run_idx]

                # Apply the translation if the target run was found
                if target_run:
                    # Verify original text match as a safety check? Maybe too strict.
                    # if target_run.text == original_text:
                    target_run.text = translated_text
                    applied_count += 1
                    # else:
                    #     log_message(f"Text mismatch for run {detailed_shape_id}. Expected '{original_text}', found '{target_run.text}'. Skipping.", level="WARNING")
                    #     skipped_count += 1
                else:
                     log_message(f"Could not find target run for entry: {entry['shape_id']}", level="WARNING")
                     skipped_count += 1

            except IndexError:
                 log_message(f"Index error applying translation for {entry['shape_id']}", level="WARNING")
                 skipped_count += 1
            except Exception as e:
                log_message(f"Error applying translation for {entry['shape_id']}: {e}", level="ERROR")
                skipped_count += 1
        
        log_message(f"Applied {applied_count} translations, skipped {skipped_count}.", level="INFO")
        prs.save(output_path)
        log_message(f"Translated presentation saved to: {output_path}", level="SUCCESS")

    except Exception as e:
        log_message(f"Error during translation process: {e}", level="ERROR")


def main():
    """Main function to handle command-line arguments and process presentations."""

    parser = argparse.ArgumentParser(description="Translate PowerPoint presentations using the OpenAI API.")
    parser.add_argument("input_path", nargs='?', type=str, help="Path to the input PowerPoint file or directory.")
    parser.add_argument("-o", "--output", type=str, help="Output file or directory path.  If not specified, defaults to [original_filename]_translated.[ext]")
    parser.add_argument("-l", "--language", type=str, required=True, help="Target language code (e.g., es, fr, zh-CN)")
    parser.add_argument("-k", "--api_key", type=str, help="OpenAI API Key. If not provided, will check the OPENAI_API_KEY environment variable.")
    parser.add_argument("-p", "--profile", action="store_true", help="Enable profiling.")

    args = parser.parse_args()

    if args.api_key:
        os.environ["OPENAI_API_KEY"] = args.api_key
    elif not os.environ.get("OPENAI_API_KEY"):
        log_message("OpenAI API Key not found. Please set the OPENAI_API_KEY environment variable or use the -k option.", level="ERROR")
        return

    input_path = args.input_path
    output_path = args.output
    target_language = args.language

    if not input_path:
        log_message("No input file or directory specified.  Processing all .pptx files in the current directory.", level="INFO")
        for filename in os.listdir("."):
            if filename.lower().endswith(SUPPORTED_EXTENSIONS):
                default_output_path = filename.replace(".pptx", f"_translated_{target_language}.pptx")
                asyncio.run(translate_presentation(filename, target_language, default_output_path))
        return

    if os.path.isfile(input_path):
        if not output_path:
            base, ext = os.path.splitext(input_path)
            output_path = f"{base}_translated_{target_language}{ext}"
        if args.profile:
            cProfile.run(f"translate_presentation('{input_path}', '{target_language}', '{output_path}')", sort="tottime")
        else:
            asyncio.run(translate_presentation(input_path, target_language, output_path))
    elif os.path.isdir(input_path):
        if not output_path:
            output_path = input_path
        elif not os.path.isdir(output_path):
            log_message("If input is a directory, output must also be a directory (or not specified).", level="ERROR")
            return

        for filename in os.listdir(input_path):
            if filename.lower().endswith(SUPPORTED_EXTENSIONS):
                full_input_path = os.path.join(input_path, filename)
                base, ext = os.path.splitext(filename)
                full_output_path = os.path.join(output_path, f"{base}_translated_{target_language}{ext}")
                if args.profile:
                    cProfile.run(f"translate_presentation('{full_input_path}', '{target_language}', '{full_output_path}')", sort="tottime")
                else:
                    asyncio.run(translate_presentation(full_input_path, target_language, full_output_path))
    else:
        log_message(f"Invalid input path: {input_path}", level="ERROR")


if __name__ == "__main__":
    main()
